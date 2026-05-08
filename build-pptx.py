"""Build Frontier Canvas GTM Presentation — Executive Deck Orchestrator output.
Winston narrative (empowerment + contributions) + consulting layouts.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
MS_BLUE = RGBColor(0x00, 0x78, 0xD4)
MS_DARK = RGBColor(0x24, 0x29, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT = RGBColor(0x50, 0xE6, 0xFF)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
DARK_GRAY = RGBColor(0x55, 0x55, 0x55)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def dark_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = MS_DARK
    return slide


def light_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    return slide


def txt(slide, left, top, width, height, text, size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════════
s = dark_slide()
txt(s, Inches(1), Inches(2), Inches(11), Inches(1.5),
    "Frontier Canvas", size=54, bold=True, color=WHITE)
txt(s, Inches(1), Inches(3.5), Inches(11), Inches(1),
    "One AE + Copilot CLI = Production App From Scratch",
    size=28, color=ACCENT)
txt(s, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
    "Edison Borges  \u00b7  ATU Americas  \u00b7  May 2026",
    size=16, color=LIGHT_GRAY)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.3), Inches(3), Pt(4))
bar.fill.solid()
bar.fill.fore_color.rgb = MS_BLUE
bar.line.fill.background()

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: ABOUT ME — Edison Borges (personal intro)
# ═══════════════════════════════════════════════════════════════
s = dark_slide()

# Photo placeholder (left side)
photo_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(1), Inches(1.5), Inches(3), Inches(3.5))
photo_box.fill.solid()
photo_box.fill.fore_color.rgb = RGBColor(0x3A, 0x3F, 0x47)
photo_box.line.color.rgb = MS_BLUE
photo_box.line.width = Pt(2)
txt(s, Inches(1.2), Inches(2.8), Inches(2.6), Inches(0.8),
    "\U0001F4F7  Add your photo here", size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Name & title (right side)
txt(s, Inches(4.8), Inches(1.3), Inches(8), Inches(1),
    "Edison Borges", size=44, bold=True, color=WHITE)
txt(s, Inches(4.8), Inches(2.3), Inches(8), Inches(0.6),
    "Account Technology Strategist  \u00b7  ATU Americas", size=20, color=ACCENT)

# Divider line
div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.8), Inches(3.1), Inches(6), Pt(2))
div.fill.solid()
div.fill.fore_color.rgb = MS_BLUE
div.line.fill.background()

# Personal details — richer bio
details = [
    ("\U0001F30E", "Based in Americas \u2014 working across industries and enterprise accounts"),
    ("\U0001F3AF", "Focused on Frontier Transformation, AI adoption & customer value realization"),
    ("\U0001F6E0\uFE0F", "Built Frontier Canvas solo using Copilot CLI \u2014 the same tools we sell"),
    ("\U0001F4AC", "Passionate about turning field insights into scalable tools that help the entire ATU"),
    ("\U0001F91D", "Previously: [Add your background — industry, prior roles, certifications]"),
]
top = Inches(3.5)
for icon, detail in details:
    txt(s, Inches(4.8), top, Inches(8), Inches(0.6),
        f"{icon}  {detail}", size=16, color=WHITE)
    top += Inches(0.7)

note(s, "Personalize this slide: add your photo, update background line, keep it under 30 seconds. The credibility IS the tool you built.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 3: PROBLEM
# ═══════════════════════════════════════════════════════════════
s = light_slide()
txt(s, Inches(0.8), Inches(0.5), Inches(11.5), Inches(1.2),
    "AEs Walk Into CxO Meetings Cold \u2014 That\u2019s a Revenue Leak",
    size=36, bold=True, color=MS_DARK)

pillars = [
    ("\U0001F6AB", "No structured prep tool", "for C-level AI conversations"),
    ("\U0001F4CA", "40-slide decks don\u2019t fit", "a 20-minute CxO slot"),
    ("\U0001F504", "Every AE reinvents the wheel", "no evidence-to-insight flow"),
]
for i, (icon, line1, line2) in enumerate(pillars):
    left = Inches(1 + i * 4)
    txt(s, left, Inches(2.5), Inches(3.5), Inches(1), icon, size=48, color=MS_BLUE, align=PP_ALIGN.CENTER)
    txt(s, left, Inches(3.5), Inches(3.5), Inches(0.8), line1, size=20, bold=True, color=MS_DARK, align=PP_ALIGN.CENTER)
    txt(s, left, Inches(4.2), Inches(3.5), Inches(0.8), line2, size=16, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4: SOLUTION — expanded with email, guidance, AE value
# ═══════════════════════════════════════════════════════════════
s = light_slide()
txt(s, Inches(0.8), Inches(0.3), Inches(11.5), Inches(1),
    "10 Minutes \u2192 8 Customer-Ready Deliverables",
    size=36, bold=True, color=MS_DARK)

principles = [
    ("\U0001F916", "Copilot-Native Engine",
     "Chat + Cowork powers all outputs\nNo custom AI backend needed"),
    ("\U0001F4E7", "Email Communications",
     "Executive briefing & follow-up emails\nready to send from Outlook"),
    ("\U0001F4CB", "Step-by-Step AE Guidance",
     "Conversation guide with talking points,\nobjection handling & next steps"),
    ("\U0001F3AF", "Customer Presentations",
     "Exec summary deck, stakeholder map,\ncompetitive differentiation"),
    ("\U0001F91D", "Account Team Handoff",
     "Complete context package for\nSA, CSM, and partner teams"),
    ("\u2705", "CELA Compliance Built In",
     "Every output follows Microsoft\nguardrails by default"),
]
for i, (icon, title, desc) in enumerate(principles):
    col = i % 3
    row = i // 3
    left = Inches(0.6 + col * 4.2)
    top = Inches(1.6 + row * 2.8)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.8), Inches(2.3))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.fill.background()
    txt(s, left + Inches(0.25), top + Inches(0.15), Inches(3.3), Inches(0.7),
        f"{icon}  {title}", size=18, bold=True, color=MS_DARK)
    txt(s, left + Inches(0.25), top + Inches(0.85), Inches(3.3), Inches(1.2),
        desc, size=14, color=DARK_GRAY)
note(s, "Key message: it's not just a dashboard — it produces ready-to-use deliverables that cover the entire meeting lifecycle: before (prep), during (conversation guide), and after (emails, handoff).")

# ═══════════════════════════════════════════════════════════════
# SLIDE 5: DEV LOOP (CLI screenshot)
# ═══════════════════════════════════════════════════════════════
s = dark_slide()
txt(s, Inches(0.8), Inches(0.3), Inches(11.5), Inches(1),
    "Copilot CLI Turns Natural Language Into Production Code",
    size=32, bold=True, color=WHITE)
if os.path.exists("cli-screenshot-3.png"):
    s.shapes.add_picture("cli-screenshot-3.png", Inches(1.5), Inches(1.5), Inches(10), Inches(5.5))
note(s, "I describe problems in natural language. Copilot plans. I review and approve. It builds, I validate, we deploy.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 6: VELOCITY KPIs — big numbers only, no crunched screenshot
# ═══════════════════════════════════════════════════════════════
s = dark_slide()
txt(s, Inches(0.8), Inches(0.5), Inches(11.5), Inches(1),
    "One Person \u00b7 Zero Cost \u00b7 Production Grade",
    size=36, bold=True, color=WHITE)

kpis = [
    ("120", "commits"),
    ("22K+", "lines of code"),
    ("85", "customer stories"),
    ("101", "use cases"),
]
for i, (num, label) in enumerate(kpis):
    left = Inches(0.8 + i * 3.1)
    txt(s, left, Inches(2.2), Inches(3), Inches(2),
        num, size=72, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txt(s, left, Inches(4.2), Inches(3), Inches(0.8),
        label, size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

txt(s, Inches(1), Inches(5.8), Inches(11), Inches(1),
    "Built entirely with GitHub Copilot CLI  \u00b7  React + TypeScript + Vite  \u00b7  Hosted on GitHub Pages at $0",
    size=16, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7: PRODUCT DEMO (2x2 grid)
# ═══════════════════════════════════════════════════════════════
s = light_slide()
txt(s, Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.8),
    "Profile \u2192 Challenges \u2192 Summary \u2192 Action Center",
    size=28, bold=True, color=MS_DARK, align=PP_ALIGN.CENTER)

screenshots = [
    ("demo-screenshots/01-profile.png", "Step 1: Customer Profile"),
    ("demo-screenshots/02-challenges.png", "Step 2: Industry Challenges"),
    ("demo-screenshots/03-exec-summary.png", "Step 3: Executive Summary"),
    ("demo-screenshots/04-action-center.png", "Step 4: Action Center"),
]
for i, (path, label) in enumerate(screenshots):
    col = i % 2
    row = i // 2
    left = Inches(0.5 + col * 6.4)
    top = Inches(1.2 + row * 3.1)
    if os.path.exists(path):
        s.shapes.add_picture(path, left, top, Inches(6), Inches(2.8))
    txt(s, left, top + Inches(2.85), Inches(6), Inches(0.4),
        label, size=14, color=MID_GRAY, align=PP_ALIGN.CENTER)
note(s, "Profile captures context. Challenges maps priorities. Summary builds the story. Action Center generates 8 customer-ready deliverables.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 8: ITERATION — 9 feedbacks (3x), two-column layout
# ═══════════════════════════════════════════════════════════════
s = light_slide()
txt(s, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
    "Real Field Feedback Ships Same Day \u2014 Not Next Quarter",
    size=32, bold=True, color=MS_DARK)

txt(s, Inches(0.6), Inches(1.3), Inches(5.5), Inches(0.4), "FEEDBACK", size=12, bold=True, color=MS_BLUE)
txt(s, Inches(6.3), Inches(1.3), Inches(0.5), Inches(0.4), "", size=12)
txt(s, Inches(7), Inches(1.3), Inches(5.5), Inches(0.4), "SHIPPED", size=12, bold=True, color=MS_BLUE)

feedbacks = [
    ("\u201CSmart Fill isn\u2019t extracting stakeholders\u201D",
     "Broadened to all C-level/VP+ roles"),
    ("\u201CToo many prompts for one meeting\u201D",
     "Merged into 1 Customer Presentation Package"),
    ("\u201CWhich Copilot do I use for what?\u201D",
     "Distinct Chat vs Cowork buttons with badges"),
    ("\u201CCompetitive content only maps to M365 Copilot\u201D",
     "Expanded to Azure AI, Foundry, Fabric, PBI"),
    ("\u201CLinks in Cowork output aren\u2019t clickable\u201D",
     "Full URLs with proper markdown formatting"),
    ("\u201CSponsors shown as confirmed — customer didn\u2019t agree\u201D",
     "Changed to \u2018Suggested Sponsors\u2019 everywhere"),
    ("\u201CUse cases say TOP — feels presumptuous\u201D",
     "Renamed to \u2018Potential Use Cases\u2019 across all prompts"),
    ("\u201CTimeline defaults to 60 days — not always realistic\u201D",
     "Cowork now asks the user for their timeline"),
    ("\u201CExec summary doesn\u2019t use customer brand\u201D",
     "Added full brand identity extraction (colors+fonts+logo)"),
]
for i, (fb, shipped) in enumerate(feedbacks):
    top = Inches(1.75 + i * 0.6)
    txt(s, Inches(0.6), top, Inches(5.5), Inches(0.55),
        fb, size=13, color=MID_GRAY)
    txt(s, Inches(6.15), top, Inches(0.5), Inches(0.55),
        "\u2192", size=16, bold=True, color=MS_BLUE, align=PP_ALIGN.CENTER)
    txt(s, Inches(7), top, Inches(5.5), Inches(0.55),
        shipped, size=13, bold=True, color=MS_DARK)
note(s, "Every single one of these shipped within 24 hours of the feedback. That's the power of Copilot CLI — the dev loop is instant.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 9: CUSTOMER VOICE — template for real feedback
# ═══════════════════════════════════════════════════════════════
s = light_slide()
txt(s, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
    "What the Field Is Saying", size=36, bold=True, color=MS_DARK)
txt(s, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.5),
    "Real feedback from AEs and leaders who have used Frontier Canvas",
    size=16, color=MID_GRAY)

# 3 quote cards — first one has placeholder text, rest are blank templates
quotes = [
    ("Paola Vergara", "1:1 Feedback \u2014 May 2026",
     "\u201C[Add Paola\u2019s feedback here]\u201D"),
    ("[Name]", "[Context]",
     "\u201C[Add quote here]\u201D"),
    ("[Name]", "[Context]",
     "\u201C[Add quote here]\u201D"),
]
for i, (name, context, quote) in enumerate(quotes):
    left = Inches(0.6 + i * 4.2)
    top = Inches(2)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              left, top, Inches(3.8), Inches(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.fill.background()
    # Quote mark
    txt(s, left + Inches(0.3), top + Inches(0.2), Inches(3.2), Inches(1),
        "\u201C", size=64, bold=True, color=MS_BLUE)
    # Quote text
    txt(s, left + Inches(0.3), top + Inches(1.2), Inches(3.2), Inches(2),
        quote, size=16, color=MS_DARK)
    # Divider
    qdiv = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               left + Inches(0.3), top + Inches(3.4), Inches(3.2), Pt(1))
    qdiv.fill.solid()
    qdiv.fill.fore_color.rgb = MS_BLUE
    qdiv.line.fill.background()
    # Attribution
    txt(s, left + Inches(0.3), top + Inches(3.6), Inches(3.2), Inches(0.4),
        name, size=14, bold=True, color=MS_DARK)
    txt(s, left + Inches(0.3), top + Inches(3.95), Inches(3.2), Inches(0.4),
        context, size=12, color=MID_GRAY)
note(s, "Fill in with real quotes. Paola Vergara's feedback from today's 1:1. Add more as they come in.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 10: SCALABILITY — cross-SA, already built in
# ═══════════════════════════════════════════════════════════════
s = light_slide()
txt(s, Inches(0.8), Inches(0.5), Inches(11.5), Inches(1),
    "One Tool, Every Solution Area \u2014 Already Built In",
    size=36, bold=True, color=MS_DARK)

pillars2 = [
    ("\U0001F30D", "Cross-Solution by Design",
     "101 use cases spanning M365 Copilot,\nAzure AI, Foundry, Fabric, PBI,\nCopilot Studio & Security"),
    ("\U0001F4CA", "Evidence-Backed for Any Vertical",
     "85 customer stories + industry\nbenchmarks mapped to all\nMicrosoft AI platforms"),
    ("\U0001F512", "Security & Observability Layer",
     "Every deliverable includes\ngovernance, compliance &\ntrust by default"),
]
for i, (icon, title, desc) in enumerate(pillars2):
    left = Inches(0.8 + i * 4.2)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              left, Inches(2), Inches(3.8), Inches(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.fill.background()
    txt(s, left + Inches(0.3), Inches(2.3), Inches(3.2), Inches(1),
        icon, size=48, color=MS_BLUE, align=PP_ALIGN.CENTER)
    txt(s, left + Inches(0.3), Inches(3.3), Inches(3.2), Inches(0.7),
        title, size=20, bold=True, color=MS_DARK, align=PP_ALIGN.CENTER)
    txt(s, left + Inches(0.3), Inches(4.2), Inches(3.2), Inches(2),
        desc, size=15, color=DARK_GRAY, align=PP_ALIGN.CENTER)
note(s, "Key message: this is NOT a tool where each SA builds their own version. It already covers all Solution Areas. The data layer (use cases, stories, challenges) spans the entire Microsoft AI portfolio.")

# ═══════════════════════════════════════════════════════════════
# SLIDE 11: CONTRIBUTIONS CLOSE
# ═══════════════════════════════════════════════════════════════
s = dark_slide()
txt(s, Inches(0.8), Inches(0.8), Inches(11.5), Inches(1.2),
    "What You Now Have That You Didn\u2019t 12 Minutes Ago",
    size=36, bold=True, color=WHITE)

contributions = [
    "A production tool that turns 10 min of prep into 8 customer-ready deliverables",
    "Living proof that the tools we sell accelerate US \u2014 not just customers",
    "A cross-Solution-Area engine backed by 85 stories and 101 use cases",
]
for i, contrib in enumerate(contributions):
    top = Inches(2.8 + i * 1.4)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), top + Inches(0.12), Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = MS_BLUE
    dot.line.fill.background()
    txt(s, Inches(1.7), top, Inches(10), Inches(1), contrib, size=24, color=WHITE)

txt(s, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
    "This slide stays up during Q&A.", size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12: APPENDIX
# ═══════════════════════════════════════════════════════════════
s = light_slide()
txt(s, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
    "Appendix: The 8 Cowork Prompts", size=28, bold=True, color=MS_DARK)

prompts = [
    ("\U0001F4C4", "Executive Summary Deck", "HTML + PPTX"),
    ("\U0001F4AC", "Conversation Guide", "HTML + Word"),
    ("\U0001F3AF", "Customer Presentation", "HTML + PPTX"),
    ("\U0001F91D", "Stakeholder Alignment", "HTML"),
    ("\U0001F4E7", "Executive Briefing Email", "Outlook draft"),
    ("\U0001F4E8", "Follow-Up Email", "Outlook draft"),
    ("\U0001F4CB", "Account Team Handoff", "HTML"),
    ("\U0001F3C6", "Competitive Differentiation", "HTML"),
]
for i, (icon, name, output) in enumerate(prompts):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 6.3)
    top = Inches(1.3 + row * 1.4)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.8), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.fill.background()
    txt(s, left + Inches(0.2), top + Inches(0.1), Inches(4.5), Inches(0.7),
        f"{icon}  {name}", size=18, bold=True, color=MS_DARK)
    txt(s, left + Inches(0.2), top + Inches(0.7), Inches(4.5), Inches(0.5),
        f"Output: {output}", size=14, color=MID_GRAY)

# Save
output_path = "Frontier-Canvas-GTM-Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
